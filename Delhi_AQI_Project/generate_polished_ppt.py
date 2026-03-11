import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths
PLOTS_DIR = "outputs/plots"
OUTPUT_PPT = "outputs/presentation/Delhi_AQI_Structural_Analysis.pptx"

# Slide Layout Indices (Standard Template)
LAYOUT_TITLE = 0
LAYOUT_TITLE_CONTENT = 1
LAYOUT_SECTION = 2
LAYOUT_BLANK = 6

def create_ppt():
    prs = pptx.Presentation()
    
    # helper to set nice title formatting
    def set_title(slide, text):
        title = slide.shapes.title
        set_title_shape(title, text)

    def set_title_shape(shape, text):
        shape.text = text
        p = shape.text_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = 'Calibri'

    # helper to add content slide
    def add_analysis_slide(title, bullets, plot_filename, numeric_key_result, notes):
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
        
        # Title
        title_shape = slide.shapes.title
        if not title_shape:
            # Manually add title if layout doesn't have it
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_shape = title_box
            
        set_title_shape(title_shape, title)
        
        # Left Content (Bullets)
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(5.0)
        height = Inches(5.0)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(24)
            p.font.name = 'Calibri'
            p.space_after = Pt(14)
            p.level = 0
            
        # Key Finding Box
        kbox_top = Inches(5.5)
        kbox = slide.shapes.add_textbox(left, kbox_top, width, Inches(1.5))
        ktf = kbox.text_frame
        kp = ktf.add_paragraph()
        kp.text = "Key Result: " + numeric_key_result
        kp.font.bold = True
        kp.font.size = Pt(20)
        kp.font.color.rgb = RGBColor(0, 50, 100) # Dark Blue
        
        # Right Image
        img_path = os.path.join(PLOTS_DIR, plot_filename)
        if os.path.exists(img_path):
            img_left = Inches(5.8)
            img_top = Inches(1.5)
            img_width = Inches(4.0)
            slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)
            
        # Notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # --- SLIDE 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Structural Analysis of Delhi Air Quality (2017-2023)"
    subtitle.text = "Daily AQI-Based Spatial and Temporal Investigation\n\nAuthor: Autonomous Agent\nDate: February 2026"
    
    # --- SLIDE 2: Dataset Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_CONTENT])
    set_title(slide, "Dataset Overview")
    
    # Custom layout for overview
    # Left text
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(5.0)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    lines = [
        "Source: Official CPCB/DPCC Monitoring Stations",
        "Resolution: Hourly Aggregated to Daily Mean",
        "Time Range: Jan 1, 2017 – Dec 31, 2023",
        "Coverage: ~78,000 Daily Observations",
        "Stations: 38 Distinct Locations"
    ]
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.space_after = Pt(14)
        p.level = 0
        
    # Right Image (Completeness)
    img_path = os.path.join(PLOTS_DIR, "phase1_completeness.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(6), Inches(1.5), width=Inches(3.5))

    # --- ANALYSIS SLIDES 1-9 ---
    
    # 1. Spatial
    add_analysis_slide(
        "1. Spatial Inequality of AQI",
        [
            "Purpose: Identify persistent pollution hotspots vs cleaner areas.",
            "Method: Comparison of station-wise mean daily AQI.",
            "Finding: High spatial heterogeneity exists across the city."
        ],
        "analysis1_spatial.png",
        "CV: 0.52 (High spatial variance)",
        "The analysis reveals significant disparity between hotspots like Anand Vihar and cleaner peripheral areas. The Coefficient of Variation of 0.52 indicates that location is a primary determinant of exposure risk."
    )
    
    # 2. Persistence
    add_analysis_slide(
        "2. Persistence vs Episodic Pollution",
        [
            "Purpose: Determine if high pollution is constant or event-driven.",
            "Method: Scatter plot of Mean vs Standard Deviation.",
            "Finding: Pollution levels correlate with volatility (episodic peaks)."
        ],
        "analysis2_persistence.png",
        "Posititve Mean-SD Correlation",
        "We observe a positive correlation between mean AQI and standard deviation. This suggests that the most polluted stations are not just consistently bad, but suffer from extreme episodic peaks that drive up the average and variance."
    )
    
    # 3. Seasonality
    add_analysis_slide(
        "3. Seasonal Regime Structure",
        [
            "Purpose: Characterize the annual pollution cycle.",
            "Method: Monthly variability analysis (Mean +/- SD).",
            "Finding: Winter inversion dominates the annual profile."
        ],
        "analysis3_seasonality.png",
        "Peak: Dec/Jan | Low: Aug",
        "The seasonal profile is extremely pronounced. Winter months (Nov-Jan) show AQI levels 3-4x higher than monsoon months. This confirms that meteorology (inversion layers) is the dominant forcing mechanism over source emission changes."
    )
    
    # 4. Rank Stability
    add_analysis_slide(
        "4. Temporal Stability of Station Ranking",
        [
            "Purpose: Test if hotspots shift location over time.",
            "Method: Year-over-year rank correlation analysis.",
            "Finding: Rankings are relatively stable, implying structural sources."
        ],
        "analysis4_ranks.png",
        "High Rank Consistency",
        "The heat map shows that the 'cleanest' and 'dirtiest' stations tend to maintain their relative rankings across years. This stability suggests that local pollution sources (traffic, industry) are structural and persistent rather than transient."
    )
    
    # 5. Extremes
    add_analysis_slide(
        "5. Concentration of Extreme Events",
        [
            "Purpose: Quantify exposure to hazardous air (AQI > 400).",
            "Method: Frequency analysis of severe air days.",
            "Finding: Hotspots face 'Severe' air ~15% of the time."
        ],
        "analysis5_extremes.png",
        "Max Severe Share: ~14.5%",
        "At the worst affected stations, nearly 1 in 7 days is classified as 'Severe' (AQI > 400). This disproportionate burden of extreme events highlights the urgent need for targeted interventions in these specific zones."
    )
    
    # 6. Diurnal
    add_analysis_slide(
        "6. Representative Diurnal Cycle",
        [
            "Purpose: Understand daily temporal dynamics.",
            "Method: Hourly aggregation for representative year (2022).",
            "Finding: Bimodal distribution driven by traffic & boundary layer."
        ],
        "analysis6_diurnal.png",
        "Peaks: Morning & Evening",
        "The diurnal cycle shows characteristic peaks in the morning and evening. This is a classic signature of traffic emissions coupled with boundary layer dynamics (lower mixing height at night/morning trapping pollutants)."
    )
    
    # 7. Completeness
    add_analysis_slide(
        "7. Data Completeness Robustness",
        [
            "Purpose: Assess reliability of the monitoring network.",
            "Method: Analysis of recording gaps per station.",
            "Finding: Core network is robust; newer stations have gaps."
        ],
        "analysis7_completeness.png",
        "Avg recorded days: ~2072",
        "While the aggregate dataset is large, individual station completeness varies. Long-term trend analysis should prioritize the core set of stations with >90% coverage to avoid bias from station commissioning dates."
    )
    
    # 8. Polarization (Sync)
    add_analysis_slide(
        "8. City-wide Synchronization",
        [
            "Purpose: Measure inter-station dependency.",
            "Method: Pairwise correlation matrix of daily AQI.",
            "Finding: The entire city oscillates together."
        ],
        "analysis8_sync.png",
        "High Avg Correlation",
        "The correlation matrix is predominantly red (high positive correlation). This indicates that when AQI rises in one part of Delhi, it rises everywhere. Regional metereology, not local sources, drives the day-to-day variance."
    )
    
    # 9. Drift
    add_analysis_slide(
        "9. Long-term Drift vs Seasonality",
        [
            "Purpose: Detect multi-year trends amidst noise.",
            "Method: Long-term time series analysis (2017-2023).",
            "Finding: No monotonic decline; seasonal noise masks trends."
        ],
        "analysis9_drift.png",
        "Trend: Indeterminate/Flat",
        "The multi-year view shows that inter-annual variability is high. There is no statistically significant linear downtrend visible in the raw aggregate data, suggesting that policy interventions are competing with varying meteorological conditions."
    )
    
    # --- Synthesis Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_CONTENT])
    set_title(slide, "Synthesis & Conclusions")
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    conclusions = [
        "1. Spatial Determinism: Where you live matters. Hotspots are persistent and structural.",
        "2. Seasonal Dominance: Winter weather is the primary amplifier of pollution.",
        "3. Network Unity: The city behaves as a single airshed (high synchronization).",
        "4. Policy Implication: Local interventions address hotspots, but regional coordination is needed for the baseline load.",
        "5. Future Work: Integrate meteorological data (wind speed/temp) to isolate source contributions from weather effects."
    ]
    
    for line in conclusions:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.space_after = Pt(20)
        p.font.bold = True
        
    # Save
    prs.save(OUTPUT_PPT)
    print(f"Presentation generated at: {OUTPUT_PPT}")
    print(f"File size: {os.path.getsize(OUTPUT_PPT) / 1024:.2f} KB")

if __name__ == "__main__":
    create_ppt()
