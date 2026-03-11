import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from zipfile import ZipFile

# Configuration
OUTPUT_REL_PATH = "outputs/presentation/Delhi_AQI_Structural_Analysis.pptx"
OUTPUT_PATH = os.path.abspath(OUTPUT_REL_PATH)
PLOTS_DIR = os.path.abspath("outputs/plots")

# Layouts (Standard)
LAYOUT_TITLE = 0
LAYOUT_TITLE_CONTENT = 1
LAYOUT_BLANK = 6

def create_ppt():
    print("Initializing Presentation...")
    prs = pptx.Presentation()

    # --- HELPER FUNCTIONS ---
    def add_slide(title_text, bullets, plot_filename=None, layout_idx=LAYOUT_BLANK):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Title
        title_shape = slide.shapes.title
        if not title_shape:
            # Manually create title if layout doesn't have one
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        
        title_shape.text = title_text
        try:
            p = title_shape.text_frame.paragraphs[0]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.name = 'Arial'
        except: pass

        # Bullet Content
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(5.0) if plot_filename else Inches(9.0)
        height = Inches(5.0)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(24)
            p.font.name = 'Arial'
            p.space_after = Pt(14)
            p.level = 0

        # Plot
        if plot_filename:
            full_plot_path = os.path.join(PLOTS_DIR, plot_filename)
            if os.path.exists(full_plot_path):
                slide.shapes.add_picture(full_plot_path, Inches(5.8), Inches(1.5), width=Inches(4.0))

    # --- SLIDE GENERATION ---

    # 1. Title Slide
    add_slide("Delhi AQI Structural Analysis (2017-2023)", 
              ["A Comprehensive Spatial & Temporal Investigation", "Author: Autonomous Analysis Agent", "Date: 2026"], 
              layout_idx=LAYOUT_TITLE_CONTENT)

    # 2. Data & Methodology
    add_slide("Data & Methodology Overview", 
              ["Source: Official CPCB/DPCC Monitoring Network", 
               "Data: ~78,000 Daily Aggregated Observations", 
               "Coverage: 38 Stations (2017-2023)",
               "Methods: Statistical Distributions, Rank Stability, Mann-Kendall Trend"],
              "phase1_completeness.png")

    # 3. Spatial Inequality
    add_slide("Spatial Inequality of AQI", 
              ["Significant pollution hotspots exist (e.g., Anand Vihar).", 
               "Peripheral areas are consistently cleaner.", 
               "High Coefficient of Variation (0.52) confirms inequality.", 
               "Location is a primary determinant of exposure."],
              "analysis1_spatial.png")

    # 4. Seasonal Winter Dominance
    add_slide("Seasonal Winter Dominance", 
              ["Winter AQI is 3-4x higher than Monsoon AQI.", 
               "Inversion layers trap pollutants in Dec/Jan.", 
               "Meteorology drives the primary annual variance.", 
               "Summer heatwave dust events cause secondary peaks."],
              "analysis3_seasonality.png")

    # 5. Inter-station Synchronization
    add_slide("Inter-station Synchronization", 
              ["High city-wide correlation (Mean r = 0.90).", 
               "The entire city oscillates as a single unit.", 
               "Regional meteorology overrides local source variance.", 
               "Policy implication: Regional coordination is essential."],
              "analysis8_sync.png")

    # 6. Long-term Trend Assessment
    add_slide("Long-term Trend Assessment", 
              ["Mann-Kendall Test detects 'No Significant Trend' (p=0.74).", 
               "Inter-annual meteorological variability masks improvements.", 
               "No linear decline visible in raw aggregate data.", 
               "Policy efforts compete with adverse weather extrema."],
              "analysis9_drift.png")

    # 7. Structural Interpretation Summary
    add_slide("Structural Interpretation Summary", 
              ["1. Spatial Determinism: Hotspots are structural and persistent.", 
               "2. Seasonal Dominance: Weather amplifies pollution 4x.", 
               "3. Network Unity: Delhi acts as one synchronized airshed.", 
               "4. Volatility: Dirtier stations are also more episodic."], 
              "analysis2_persistence.png")

    # 8. Limitations
    add_slide("Limitations of Analysis", 
              ["• Composite AQI masks individual pollutant trends (PM2.5 vs Ozone).", 
               "• Meteorology not explicitly modeled (Wind/Temp).", 
               "• Spatial gaps in Western/Northern periphery.", 
               "• Correlation does not prove causation."])

    # 9. Next Research Directions
    add_slide("Next Research Directions", 
              ["• Pollutant-specific trend analysis (PM2.5 only).", 
               "• Meteorology-adjustment (De-weathering) techniques.", 
               "• Localized exposure mapping using interpolation.", 
               "• Early warning forecasting models."])

    # 10. Final Conclusion
    add_slide("Final Conclusion", 
              ["Delhi's pollution is structurally rigid and seasonally amplified.", 
               "While hotspots require local enforcement, the baseline load is regional.", 
               "Future gains require attacking the regional baseline.", 
               "Current trends show stagnation, necessitating aggressive policy shifts."])

    # --- SAVE ---
    print(f"Saving to {OUTPUT_PATH}...")
    if os.path.exists(OUTPUT_PATH):
        os.remove(OUTPUT_PATH)
        
    prs.save(OUTPUT_PATH)
    print("Save complete.")
    return len(prs.slides)

def validate_ppt():
    print(f"Validating {OUTPUT_PATH}...")
    try:
        if not os.path.exists(OUTPUT_PATH):
            raise FileNotFoundError("Output file not found.")
            
        with ZipFile(OUTPUT_PATH, 'r') as z:
            bad_file = z.testzip()
            if bad_file:
                raise Exception(f"Corrupt file segment: {bad_file}")
            # Check for core xml to be sure
            if 'ppt/presentation.xml' not in z.namelist():
                raise Exception("Missing ppt/presentation.xml - not a valid PPTX.")
                
        print("VALID PPTX CONFIRMED")
        return True
    except Exception as e:
        print(f"VALIDATION FAILED: {e}")
        return False

if __name__ == "__main__":
    try:
        count = create_ppt()
        if validate_ppt():
            print(f"✔ PPT successfully created")
            print(f"✔ Path: {OUTPUT_REL_PATH}")
            print(f"✔ Slide count: {count}")
            print(f"✔ Validation = PASSED")
        else:
            print("❌ PPT Creation Failed Validation")
    except Exception as e:
        print(f"❌ Script Error: {e}")
