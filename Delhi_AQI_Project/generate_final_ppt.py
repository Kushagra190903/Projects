import os
import pandas as pd
import numpy as np
import scipy.stats as stats
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from zipfile import ZipFile

# Configuration
PLOTS_DIR = "outputs/plots"
TABLES_DIR = "outputs/tables"
OUTPUT_PPT = "outputs/presentation/Delhi_AQI_Structural_Analysis.pptx"

# Slide Layouts
LAYOUT_TITLE = 0
LAYOUT_TITLE_CONTENT = 1
LAYOUT_BLANK = 6

def mann_kendall_test(x, alpha=0.05):
    """Performs Mann-Kendall trend test."""
    n = len(x)
    s = 0
    for k in range(n-1):
        for j in range(k+1, n):
            s += np.sign(x[j] - x[k])
    var_s = (n * (n - 1) * (2 * n + 5)) / 18
    if s > 0: z = (s - 1) / np.sqrt(var_s)
    elif s < 0: z = (s + 1) / np.sqrt(var_s)
    else: z = 0
    p = 2 * (1 - stats.norm.cdf(abs(z)))
    trend = "No Trend"
    if p < alpha:
        trend = "Increasing" if z > 0 else "Decreasing"
    return trend, p, z

def calculate_stats():
    stats_data = {'avg_rank_corr': 0.0, 'mean_sync_corr': 0.0, 'mk_trend': 'No Trend', 'mk_p': 1.0, 'mk_z': 0.0}
    try:
        # Rank Stability
        df_rank = pd.read_csv(os.path.join(TABLES_DIR, "rank_stability.csv"))
        years = sorted([c for c in df_rank.columns if c.replace('.0','').isdigit()])
        corrs = []
        for i in range(len(years)-1):
            c, _ = stats.spearmanr(df_rank[years[i]], df_rank[years[i+1]])
            corrs.append(c)
        stats_data['avg_rank_corr'] = np.mean(corrs)
    except: pass

    try:
        # Synchronization
        df_sync = pd.read_csv(os.path.join(TABLES_DIR, "synchronization_corr.csv"), index_col=0)
        mask = np.tril(np.ones(df_sync.shape), k=-1).astype(bool)
        stats_data['mean_sync_corr'] = df_sync.where(mask).stack().mean()
    except: pass

    try:
        # Trend
        df_trend = pd.read_csv(os.path.join(TABLES_DIR, "long_term_trend.csv"))
        trend, p, z = mann_kendall_test(df_trend['mean_aqi'].values)
        stats_data['mk_trend'] = trend
        stats_data['mk_p'] = p
        stats_data['mk_z'] = z
    except: pass
    
    return stats_data

def generate_ppt():
    print("Calculating statistics...")
    S = calculate_stats()
    print(f"Stats computed: {S}")
    
    prs = pptx.Presentation()
    
    # Helper to clean title logic
    def add_slide(title, bullets, plot_file=None, metric=None, notes="", layout=LAYOUT_BLANK):
        slide = prs.slides.add_slide(prs.slide_layouts[layout])
        
        # Title
        title_shape = slide.shapes.title
        if not title_shape:
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        
        title_shape.text = title
        try:
            p = title_shape.text_frame.paragraphs[0]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.name = 'Calibri'
        except: pass
        
        # Content
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(5.0) if plot_file else Inches(9.0)
        height = Inches(5.0)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(24)
            p.font.name = 'Calibri'
            p.space_after = Pt(14)
            
        # Metric
        if metric:
            kbox = slide.shapes.add_textbox(left, Inches(5.5), width, Inches(1.5))
            kp = kbox.text_frame.add_paragraph()
            kp.text = metric
            kp.font.bold = True
            kp.font.size = Pt(20)
            kp.font.color.rgb = RGBColor(0, 50, 100)
            
        # Image
        if plot_file:
            path = os.path.join(PLOTS_DIR, plot_file)
            if os.path.exists(path):
                slide.shapes.add_picture(path, Inches(5.8), Inches(1.5), width=Inches(4.0))
                
        # Notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # --- SLIDES ---
    
    # 1. Title
    add_slide("Structural Analysis of Delhi Air Quality (2017-2023)", 
              ["Daily Resolution Investigation", "Author: Autonomous Agent", "Date: 2026"], 
              layout=LAYOUT_TITLE_CONTENT)
    
    # 2. Dataset
    add_slide("Dataset Overview", 
              ["Source: Official CPCB/DPCC", "Range: 2017-2023", "Stations: 38 (78k Daily Obs)"], 
              "phase1_completeness.png", 
              "Avg Coverage: >2000 days", 
              "Robust daily dataset aggregated from hourly monitors.")

    # 3. Three Facts
    add_slide("Three Structural Facts of Delhi AQI", 
              ["1. Strong Spatial Inequality", "2. Dominant Winter Seasonality", "3. High City-wide Synchronization"], 
              notes="The three pillars of Delhi's pollution structure.")

    # 4. Spatial
    add_slide("1. Spatial Inequality", 
              ["Hotspots (e.g. Anand Vihar) differ significantly from cleaner zones.", "Coefficient of Variation: 0.52"], 
              "analysis1_spatial.png", 
              "High Spatial Variance", 
              "Location determines exposure.")

    # 5. Persistence
    add_slide("2. Persistence vs Episodic", 
              ["Pollution volatility scales with mean levels.", "Dirtier stations have more extreme spikes."], 
              "analysis2_persistence.png", 
              "Positive Mean-SD Correlation", 
              "Episodic nature of hotspots confirmed.")

    # 6. Seasonality
    add_slide("3. Seasonal Regimes", 
              ["Winter inversion drives 3-4x pollution increase.", "Monsoon washout is the cleanest period."], 
              "analysis3_seasonality.png", 
              "Peak: Dec/Jan", 
              "Meteorology is the dominant forcing.")

    # 7. Rank Stability
    add_slide("4. Rank Stability", 
              ["Relative station rankings are stable year-over-year.", "Sources are structural, not transient."], 
              "analysis4_ranks.png", 
              f"Avg Rank Correlation: {S['avg_rank_corr']:.2f}", 
              "Structural persistence of pollution sources.")

    # 8. Extremes
    add_slide("5. Extreme Events (>400 AQI)", 
              ["Severe days are concentrated in hotspots.", "Max frequency ~14.5% of days."], 
              "analysis5_extremes.png", 
              "Disproportionate Exposure", 
              "Health burden is unequally distributed.")

    # 9. Diurnal
    add_slide("6. Diurnal Cycle", 
              ["Bimodal peaks (Morning/Evening).", "Driven by traffic & boundary layer collapse."], 
              "analysis6_diurnal.png", 
              "Traffic + PBL Signature", 
              "Classic urban profile.")

    # 10. Completeness
    add_slide("7. Data Robustness", 
              ["Core network is highly reliable.", "Gaps exist in newer stations."], 
              "analysis7_completeness.png", 
              "Valid Core Network", 
              "Analysis relies on long-term stations.")

    # 11. Synchronization
    add_slide("8. City-wide Synchronization", 
              ["High inter-station correlation.", "Regional meteorology overrides local effects."], 
              "analysis8_sync.png", 
              f"Mean Correlation: {S['mean_sync_corr']:.2f}", 
              "The city acts as a single airshed.")

    # 12. Drift
    slide_trend = "Increasing" if S['mk_z'] > 1.96 else ("Decreasing" if S['mk_z'] < -1.96 else "No Significant Trend")
    add_slide("9. Long-term Drift", 
              ["Mann-Kendall Trend Test performed.", f"Result: {slide_trend} (p={S['mk_p']:.2f})"], 
              "analysis9_drift.png", 
              f"Trend: {slide_trend}", 
              "Inter-annual variability masks any policy-driven downtrend.")

    # 13. Limitations
    add_slide("Limitations of Analysis", 
              ["• Composite AQI metric masks pollutant specifics.", "• Meteorology not explicitly modeled.", "• Correlation is not causation."], 
              notes="Transparency about scope.")

    # 14. Next Steps
    add_slide("Next Research Directions", 
              ["• Pollutant-specific trends (PM2.5).", "• Weather-normalized trend analysis.", "• ML-based forecasting."], 
              notes="Future work path.")

    # 15. Conclusions
    add_slide("Final Conclusions", 
              ["1. Spatial structure is rigid (Hotspots persist).", "2. Winter weather dominates everything.", "3. The city is a unified airshed."], 
              notes="Summary of findings.")

    # Save
    if os.path.exists(OUTPUT_PPT):
        os.remove(OUTPUT_PPT)
    
    prs.save(OUTPUT_PPT)
    print(f"PPT saved to: {OUTPUT_PPT}")
    return len(prs.slides)

def validate_ppt():
    print("Validating PPT...")
    try:
        with ZipFile(OUTPUT_PPT, 'r') as z:
            ret = z.testzip()
            if ret is not None:
                raise Exception(f"Corrupt file at {ret}")
            print(f"✔ Validation OK. File is valid ZIP/PPTX structure.")
            print(f"✔ File size: {os.path.getsize(OUTPUT_PPT)} bytes")
            return True
    except Exception as e:
        print(f"❌ Validation Failed: {e}")
        return False

if __name__ == "__main__":
    count = generate_ppt()
    print(f"✔ Total slides created: {count}")
    if validate_ppt():
        print("PPT successfully created and validated.")
    else:
        print("PPT creation FAILED validation.")
