from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Create presentation
prs = Presentation()

# ---------- Helper function ----------
def add_bullet_slide(title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.title
    content = slide.placeholders[1]

    title_box.text = title

    tf = content.text_frame
    tf.text = bullets[0]

    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0


# ---------- Slides ----------

add_bullet_slide(
    "Delhi AQI Structural Analysis (2017–2023)",
    [
        "Station-level hourly AQI ingested into PostgreSQL",
        "Converted to daily analytical dataset",
        "Nine structural analyses performed",
        "Goal: understand spatial, temporal, and network behaviour",
    ],
)

add_bullet_slide(
    "Spatial Structure",
    [
        "Persistent AQI hotspots exist across stations",
        "Industrial & traffic zones dominate pollution",
        "Clean zones remain relatively stable",
        "Indicates structural inequality in exposure",
    ],
)

add_bullet_slide(
    "Seasonality",
    [
        "Winter shows strongest AQI peaks",
        "Meteorology drives pollution accumulation",
        "Monsoon provides natural cleansing",
        "Seasonality dominates yearly variance",
    ],
)

add_bullet_slide(
    "City-wide Synchronization",
    [
        "Stations move together during pollution events",
        "High inter-station correlation observed",
        "Suggests regional atmospheric forcing",
        "Supports need for city-level policy action",
    ],
)

add_bullet_slide(
    "Long-Term Trend",
    [
        "No strong linear AQI decline detected",
        "Inter-annual variability remains high",
        "Policy impact not yet structurally visible",
        "Requires pollutant-level investigation",
    ],
)

add_bullet_slide(
    "Limitations",
    [
        "AQI is a composite indicator",
        "Meteorology not explicitly modeled",
        "Correlation does not imply causation",
        "Station coverage may vary",
    ],
)

add_bullet_slide(
    "Next Research Directions",
    [
        "PM2.5-specific trend analysis",
        "Meteorology-adjusted modeling",
        "Forecasting & early-warning systems",
        "Zone-level exposure mapping",
    ],
)

add_bullet_slide(
    "Conclusion",
    [
        "Delhi AQI shows strong spatial inequality",
        "Winter seasonality is dominant driver",
        "City behaves as synchronized pollution system",
        "Structural improvement not yet evident",
    ],
)

# ---------- Save ----------
output_path = "Delhi_AQI_Structural_Analysis.pptx"
prs.save(output_path)

print("✅ PPT created successfully at:", os.path.abspath(output_path))
