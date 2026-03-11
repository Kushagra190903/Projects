from pptx import Presentation
from pptx.util import Inches, Pt
import os

PPT_PATH = "outputs/presentation/Delhi_AQI_Structural_Analysis.pptx"

def create_presentation():
    """Initializes a new PowerPoint presentation."""
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Delhi AQI Structural Analysis"
    subtitle.text = "Autonomous Research Report\nDaily Resolution Analysis (2017-2023)"
    
    return prs

def add_analysis_slide(prs, title, metrics_text, interpretation, plot_path, notes=""):
    """Adds a standardized analysis slide."""
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    slide.shapes.title.text = title
    
    # Content (Metrics + Interpretation)
    # We'll use a text box for metrics/interpretation on the left, plot on the right
    # Adjusting layout manually for better fit
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(5.0)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Key Findings:"
    p.font.bold = True
    p.font.size = Pt(18)
    
    for line in metrics_text.split('\n'):
        if line.strip():
            p = tf.add_paragraph()
            p.text = f"• {line.strip()}"
            p.font.size = Pt(16)
            
    p = tf.add_paragraph()
    p.text = "\nInterpretation:"
    p.font.bold = True
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = interpretation
    p.font.size = Pt(16)
    
    # Plot
    if plot_path and os.path.exists(plot_path):
        img_left = Inches(5.2)
        img_top = Inches(1.5)
        img_width = Inches(4.5)
        slide.shapes.add_picture(plot_path, img_left, img_top, width=img_width)
        
    # Notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
        
def save_presentation(prs):
    prs.save(PPT_PATH)
    return PPT_PATH
